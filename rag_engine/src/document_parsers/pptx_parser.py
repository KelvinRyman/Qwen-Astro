"""
PPTX文档解析器

支持Microsoft PowerPoint PPTX格式文档的文本提取、元数据提取和智能分块。
"""

import logging
from typing import List, Optional
from pathlib import Path
from .base_parser import DocumentParser, DocumentMetadata, DocumentChunk

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError as e:
    logging.warning(f"PPTX解析库导入失败: {e}")
    PPTX_AVAILABLE = False


class PptxParser(DocumentParser):
    """PPTX文档解析器"""
    
    def __init__(self, config=None):
        super().__init__(config)
        self.extract_notes = self.config.get('extract_notes', True)
        self.extract_slide_titles = self.config.get('extract_slide_titles', True)
        self.min_text_length = self.config.get('min_text_length', 5)
    
    def can_parse(self, file_path: str) -> bool:
        """检查是否能解析PPTX文件"""
        if not PPTX_AVAILABLE:
            return False
        
        return Path(file_path).suffix.lower() == '.pptx'
    
    def get_supported_extensions(self) -> List[str]:
        """获取支持的文件扩展名"""
        return ['.pptx']
    
    def extract_metadata(self, file_path: str) -> DocumentMetadata:
        """提取PPTX元数据"""
        metadata = DocumentMetadata(format_type='pptx')
        
        try:
            prs = Presentation(file_path)
            
            # 基本信息
            slide_count = len(prs.slides)
            metadata.page_count = slide_count
            
            # 核心属性
            if hasattr(prs, 'core_properties'):
                core_props = prs.core_properties
                metadata.title = core_props.title
                metadata.author = core_props.author
                metadata.subject = core_props.subject
                metadata.keywords = core_props.keywords
                metadata.creator = core_props.creator
                
                # 日期
                if core_props.created:
                    metadata.creation_date = core_props.created.strftime('%Y-%m-%d %H:%M:%S')
                if core_props.modified:
                    metadata.modification_date = core_props.modified.strftime('%Y-%m-%d %H:%M:%S')
            
            # 统计信息
            total_text_shapes = 0
            total_images = 0
            total_tables = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        total_text_shapes += 1
                    elif shape.shape_type == 13:  # Picture
                        total_images += 1
                    elif shape.shape_type == 19:  # Table
                        total_tables += 1
            
            metadata.extra_metadata = {
                'slide_count': slide_count,
                'text_shapes': total_text_shapes,
                'images': total_images,
                'tables': total_tables
            }
            
            # 文件大小
            metadata.file_size = Path(file_path).stat().st_size
            
        except Exception as e:
            self.logger.error(f"提取PPTX元数据失败: {e}")
        
        return metadata
    
    def extract_text(self, file_path: str) -> str:
        """提取PPTX全文"""
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = self._extract_slide_text(slide, slide_idx)
                if slide_text:
                    text_parts.append(slide_text)
            
            return self._clean_text('\n\n'.join(text_parts))
            
        except Exception as e:
            self.logger.error(f"提取PPTX文本失败: {e}")
            return ""
    
    def extract_chunks(self, file_path: str) -> List[DocumentChunk]:
        """提取PPTX分块（按幻灯片分块）"""
        chunks = []
        base_metadata = {'file_name': Path(file_path).name, 'format_type': 'pptx'}
        
        try:
            prs = Presentation(file_path)
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = self._extract_slide_text(slide, slide_idx)
                
                if slide_text:
                    # 提取幻灯片标题
                    slide_title = self._extract_slide_title(slide)
                    
                    chunk_metadata = self._create_chunk_metadata(
                        base_metadata,
                        {
                            'slide_number': slide_idx,
                            'slide_title': slide_title,
                            'chunk_type': 'slide',
                            'page_label': str(slide_idx)
                        }
                    )
                    
                    chunk = DocumentChunk(
                        text=slide_text,
                        metadata=chunk_metadata,
                        chunk_id=f"slide_{slide_idx}",
                        page_number=slide_idx,
                        section_title=slide_title,
                        chunk_type='slide'
                    )
                    chunks.append(chunk)
            
        except Exception as e:
            self.logger.error(f"提取PPTX分块失败: {e}")
        
        return chunks
    
    def _extract_slide_text(self, slide, slide_number: int) -> str:
        """提取单张幻灯片的文本"""
        text_parts = []
        
        # 提取形状中的文本
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                if len(text) >= self.min_text_length:
                    text_parts.append(text)
            
            # 提取表格文本
            elif hasattr(shape, 'table'):
                table_text = self._extract_table_text(shape.table)
                if table_text:
                    text_parts.append(table_text)
        
        # 提取备注
        if self.extract_notes and hasattr(slide, 'notes_slide'):
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text and len(notes_text) >= self.min_text_length:
                    text_parts.append(f"备注: {notes_text}")
            except Exception:
                pass
        
        slide_text = '\n\n'.join(text_parts)
        return self._clean_text(slide_text) if slide_text else ""
    
    def _extract_slide_title(self, slide) -> Optional[str]:
        """提取幻灯片标题"""
        if not self.extract_slide_titles:
            return None
        
        try:
            # 通常标题是第一个文本框或者具有特定布局的形状
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    # 简单启发式：较短的文本可能是标题
                    text = shape.text.strip()
                    if len(text) < 100 and '\n' not in text:
                        return text
            
            # 如果没有找到明显的标题，返回第一个文本
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    # 取第一行作为标题
                    first_line = text.split('\n')[0]
                    if len(first_line) < 100:
                        return first_line
                    
        except Exception:
            pass
        
        return None
    
    def _extract_table_text(self, table) -> str:
        """提取表格文本"""
        try:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_data.append(cell_text)
                if row_data:
                    table_data.append(' | '.join(row_data))
            
            return '\n'.join(table_data)
        except Exception as e:
            self.logger.warning(f"提取表格文本失败: {e}")
            return ""
